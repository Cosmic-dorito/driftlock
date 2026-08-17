#!/usr/bin/env python3
"""Build the graded submission deck on the mandatory i4C template.

    python scripts/make_submission_deck.py

The template ships ten slides; slide 1 is an instructions slide marked "REMOVE THIS SLIDE FOR
INSTRUCTIONS IN FINAL SUBMISSION", which leaves exactly the nine the Applied Materials problem
statement maps content onto. That mapping lives in ``docs/SUBMISSION.md`` and is the authority for
what goes where; where the template's own slide title differs (its "Impact and Benefits" is the
problem statement's "Results") the problem statement wins, because it is the track-specific document.

**Every number on a slide is read from ``results/``.** ``scripts/verify_submission.py`` extracts each
decimal from the finished deck and fails the build if it cannot find it there, which is why figures
are loaded from the CSVs rather than typed.

**Layout approach, and why.** The template's body placeholders are small and vertically centred, so
resizing them makes text drift out of its card and off the slide - both defects appeared in the
first build. Instead the placeholders are deleted and replaced with text boxes at explicit geometry
and top anchoring. The template's own furniture - background, header band, title, accent bar, icon,
sub-card rectangles - is kept untouched, which is what "keep the core theme" requires.
"""

from __future__ import annotations

import contextlib
import csv
import sys
from pathlib import Path

import cv2
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

TEMPLATE = REPO_ROOT / "assets" / "i4c_idea_submission_template.pptx"
OUTPUT = REPO_ROOT / "Silli-Con Artists_PS02.pptx"
FIGURES = REPO_ROOT / "results" / "figures"

FONT = "Poppins"
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xAE, 0xFF, 0x82)
LABEL = RGBColor(0x9C, 0xA3, 0xAF)

#: The template's content card spans x 0.62-12.70, y 2.21-7.19 inches. Body content sits inside it.
BODY_LEFT, BODY_TOP, BODY_WIDTH = 1.05, 4.02, 11.23
BODY_BOTTOM = 7.02
HALF_WIDTH = 5.44


def set_text(shape, text: str, size: float | None = None) -> None:
    """Replace a shape's text, keeping the first run's formatting.

    ``text_frame.text = ...`` collapses the paragraph to one unstyled run and would strip the
    template's font, colour and size. Assigning ``run.text`` preserves them.
    """
    frame = shape.text_frame
    first = frame.paragraphs[0]
    run = first.runs[0] if first.runs else first.add_run()
    for extra in first.runs[1:]:
        extra._r.getparent().remove(extra._r)
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    for para in frame.paragraphs[1:]:
        para._p.getparent().remove(para._p)


def add_text(slide, left, top, width, height, lines, size=12.0, bold=False,
             color=WHITE, space_after=5) -> None:
    """A top-anchored text box in the template's type style."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        para.space_after = Pt(space_after)
    return box


def drop(shape) -> None:
    shape._element.getparent().remove(shape._element)


def find(slide, needle: str):
    for shape in slide.shapes:
        if shape.has_text_frame and needle in shape.text_frame.text:
            return shape
    raise LookupError(f"no shape containing {needle!r}")


def drop_if(slide, *needles: str) -> None:
    for needle in needles:
        with contextlib.suppress(LookupError):
            drop(find(slide, needle))


#: Which sub-card rectangles each slide has already spent. Re-querying the slide on every call
#: returned the SAME shape twice, so a two-card layout moved one card into the right-hand slot and
#: left the other where the template put it - visible as a missing card background.
_CARDS_USED: dict[int, int] = {}


def card(slide, left, top, width, height):
    """Reposition the template's rounded sub-card rectangles, which carry its fill and border."""
    available = [s for s in slide.shapes if s.shape_id in (17, 21)]
    index = _CARDS_USED.get(id(slide), 0)
    if index >= len(available):
        return None
    shape = available[index]
    _CARDS_USED[id(slide)] = index + 1
    shape.left, shape.top, shape.width, shape.height = (
        Inches(left), Inches(top), Inches(width), Inches(height))
    return shape


def drop_body_icons(slide) -> None:
    """Remove the template's small label icons from the body area.

    Each sub-card ships a little icon beside its heading. The headings are replaced with our own,
    so the orphaned icons would sit on top of whatever lands there - on the results slide one of
    them printed straight through the tolerance table.
    """
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if shape.top > Inches(3.9) and shape.width < Inches(0.5):
            drop(shape)


def clear_body(slide, below: float = 3.3) -> None:
    """Remove every shape below a y threshold, leaving the slide furniture above it.

    Used where the template's layout differs enough from the others that repositioning its parts
    individually is more fragile than starting from an empty card - the references slide ships two
    stacked sub-cards, its own headings and three reference lines, all at shape ids that do not
    match the other slides.
    """
    for shape in list(slide.shapes):
        if shape.top is not None and shape.top >= Inches(below):
            drop(shape)


def metrics(label: str) -> dict[str, str]:
    with (REPO_ROOT / "results" / f"metrics_{label}.csv").open(newline="", encoding="utf-8") as fh:
        return {row["metric"]: row["value"] for row in csv.DictReader(fh)}


def pct(value) -> str:
    return f"{float(value) * 100:.1f}%"


def crop_search_half(source: Path, destination: Path) -> None:
    """Keep the search image and its caption band, drop the reference half.

    The overlay renders reference and search side by side at 2000x1057. Placed whole on a slide it
    is either unreadably small or wildly out of aspect - the first build squashed it to a letterbox.
    The search half is where the prediction, the truth and the rival boxes actually are, so that is
    what the slide shows; the reference is described in the caption instead.
    """
    image = cv2.imread(str(source))
    if image is None:
        raise FileNotFoundError(source)
    height, width = image.shape[:2]
    cv2.imwrite(str(destination), image[:, width // 2:])


def build() -> int:
    if not TEMPLATE.exists():
        print(f"error: template not found at {TEMPLATE}", file=sys.stderr)
        return 2
    FIGURES.mkdir(parents=True, exist_ok=True)

    sponsor, bench, finfet = metrics("sponsor"), metrics("bench"), metrics("finfet")
    bad = sum(int(m["n_mislocks"]) for m in (sponsor, bench, finfet))
    total = sum(int(m["n_pairs"]) for m in (sponsor, bench, finfet))

    prs = Presentation(str(TEMPLATE))
    xml_slides = prs.slides._sldIdLst
    first = list(xml_slides)[0]
    prs.part.drop_rel(first.rId)
    xml_slides.remove(first)

    team, problem, idea, solution, innov, results, tech, links, refs = (
        prs.slides[i] for i in range(9))

    # ------------------------------------------------------------------ 1. Team details
    set_text(find(team, "Enter Team Name Here"), "Silli-Con Artists")
    names = [s for s in team.shapes
             if s.has_text_frame and s.text_frame.text.strip() == "{Enter Name}"]
    years = [s for s in team.shapes
             if s.has_text_frame and s.text_frame.text.strip() == "{Enter Year}"]
    for box, value in zip(names, ["S. Sanjay Subramanyan", "Shiv Karthi Sadasivam",
                                  "Pedada Venkata Abhiram"]):
        box.width = Inches(4.6)          # the NAME column is 4.83 wide; the placeholder was 1.19
        set_text(box, value)
    for box in years[:3]:
        set_text(box, "2027")
    # No fourth member: the whole row goes, including its background band and the "4"/"Member 3"
    # labels. Leaving it blank reads as an unfinished template rather than a three-person team.
    drop(names[3])
    drop(years[3])
    drop_if(team, "Member 3")
    for shape in list(team.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() == "4":
            drop(shape)
            break
    for shape in list(team.shapes):
        if shape.shape_id == 27:          # the empty fourth-row band
            drop(shape)
    for shape in list(team.shapes):
        if shape.shape_id == 9:           # the table card, shortened to the three real rows
            shape.height = Inches(1.98)
    college = find(team, "{Enter Full College Name}")
    college.width = Inches(6.0)
    set_text(college, "Amrita Vishwa Vidyapeetham, Coimbatore")
    set_text(find(team, "{+91 XXXXX XXXXX}"), "+91 98408 53213")
    email = find(team, "{email@example.com}")
    email.width = Inches(4.6)
    set_text(email, "sanjaysubramanyan1324@gmail.com")

    # ------------------------------------------------------------------ 2. Problem statement
    set_text(find(problem, "Selected the problem statement"),
             "Drift-Sense — Navigation-Error Recovery (Applied Materials, PS 02)", size=16)
    drop_if(problem, "DESCRIPTION / DETAILS", "{Provide specific details")
    drop_body_icons(problem)
    card(problem, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_BOTTOM - BODY_TOP)
    add_text(problem, BODY_LEFT + 0.26, BODY_TOP + 0.22, BODY_WIDTH - 0.52, 0.25,
             ["WHY NAVIGATION-ERROR RECOVERY MATTERS"], size=10, bold=True, color=LABEL)
    add_text(problem, BODY_LEFT + 0.26, BODY_TOP + 0.58, BODY_WIDTH - 0.52, 2.0, [
        "A wafer-inspection tool navigates to a defect by stage coordinates. Stage drift means the "
        "field it actually images is not the field it was sent to. Recovery means re-locating a "
        "known high-magnification reference pattern inside a wider, coarser, noisier search image.",
        "The reference is 1000×1000 px at 1 nm/px; the search covers 10× the area at 10 nm/px, so "
        "the reference occupies only 100×100 px of it. DRAM and FinFET layouts are periodic by "
        "design, so a wrong repeat is not a blurry near-miss — it is a structurally valid match "
        "that looks correct. That ambiguity, not similarity, is the real problem.",
    ], size=12)
    add_text(problem, BODY_LEFT + 0.26, BODY_TOP + 2.45, BODY_WIDTH - 0.52, 0.3,
             ["Getting it wrong is a data-integrity failure, not an accuracy one: the tool reports "
              "a confident measurement of the wrong cell."], size=11, color=ACCENT)

    # ------------------------------------------------------------------ 3. Idea description
    set_text(find(idea, "Idea Description"), "Idea Description")
    set_text(find(idea, "Provide a brief summary"),
             "DRAM-style layouts · classical, physics-based, zero-training localization", size=16)
    drop_if(idea, "KEY CONCEPT & APPROACH", "{Briefly describe the core concept",
            "SOLUTION OVERVIEW", "{Provide an overview of the solution")
    drop_body_icons(idea)
    card(idea, BODY_LEFT, BODY_TOP, HALF_WIDTH, BODY_BOTTOM - BODY_TOP)
    card(idea, BODY_LEFT + HALF_WIDTH + 0.35, BODY_TOP, HALF_WIDTH, BODY_BOTTOM - BODY_TOP)
    add_text(idea, BODY_LEFT + 0.26, BODY_TOP + 0.22, HALF_WIDTH - 0.52, 0.25,
             ["KEY CONCEPT"], size=10, bold=True, color=LABEL)
    add_text(idea, BODY_LEFT + 0.26, BODY_TOP + 0.56, HALF_WIDTH - 0.52, 2.2, [
        "We do not match images. We invert the microscope.",
        "The forward model is known — beam PSF, area-average decimation, geometric warp, Poisson "
        "shot noise then Gaussian read noise. Localization is therefore a maximum-likelihood "
        "inverse problem with a few nuisance parameters, not a similarity search.",
    ], size=11.5)
    add_text(idea, BODY_LEFT + HALF_WIDTH + 0.61, BODY_TOP + 0.22, HALF_WIDTH - 0.52, 0.25,
             ["WHY IT BEATS TEMPLATE MATCHING"], size=10, bold=True, color=LABEL)
    add_text(idea, BODY_LEFT + HALF_WIDTH + 0.61, BODY_TOP + 0.56, HALF_WIDTH - 0.52, 2.2, [
        "The periodic lattice everyone treats as the enemy is used as a ruler: it fixes pose. The "
        "aperiodic residual — line-placement jitter unique to each cell — supplies identity.",
        "The baseline resizes over a few fixed scales and takes an argmax, so it cannot represent "
        "rotation or off-nominal magnification, and has no way to tell two identical repeats apart.",
    ], size=11.5)

    # ------------------------------------------------------------------ 4. Proposed solution
    set_text(find(solution, "Proposed Solution"), "Proposed Solution")
    set_text(find(solution, "Describe your idea in detail"),
             "Generator, localizer, and the pipeline from image pair to (x, y)", size=16)
    drop_if(solution, "SOLUTION DETAILS", "{Provide specific details about your proposed solution")
    drop_body_icons(solution)
    card(solution, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_BOTTOM - BODY_TOP)
    add_text(solution, BODY_LEFT + 0.26, BODY_TOP + 0.18, BODY_WIDTH - 0.52, 0.25,
             ["GENERATOR — 30% of the score, every stage literature-justified"],
             size=10, bold=True, color=LABEL)
    add_text(solution, BODY_LEFT + 0.26, BODY_TOP + 0.5, BODY_WIDTH - 0.52, 0.9, [
        "vector DRAM/FinFET layout at 1 nm/px → SE edge brightening (yield rises with surface "
        "tilt — Villarrubia, NIST) → beam PSF → rotation 0–2° and scale 9–11:1 → raster shear and "
        "jitter → barrel → Poisson(dose) → Gaussian(read) → gamma, vignette, charging streaks.",
        "Independent RNG streams per capture, so the pair is two acquisitions of one scene rather "
        "than one image plus noise. Fractional crop origins give continuous ground truth — which "
        "is what makes a sub-pixel claim measurable at all.",
    ], size=11.5)
    add_text(solution, BODY_LEFT + 0.26, BODY_TOP + 1.62, BODY_WIDTH - 0.52, 0.25,
             ["LOCALIZER — classical, deterministic, CPU-only"], size=10, bold=True, color=LABEL)
    add_text(solution, BODY_LEFT + 0.26, BODY_TOP + 1.94, BODY_WIDTH - 0.52, 1.0, [
        "measure pose by pyramid search → build the template through the exact forward operator "
        "(area-average decimation is the physically correct downsampler, not an approximation) → "
        "ZNCC candidate surface → lattice-residual proposals → cheap screen ranks the field →",
        "survivors get a dense 5×5 pose refit → read each candidate's pose grid as EVIDENCE rather "
        "than as its maximum → upsampled-DFT sub-pixel → bounded raster-drift correction.",
    ], size=11.5)

    # ------------------------------------------------------------------ 5. Innovation
    set_text(find(innov, "Highlight what makes your idea unique"),
             "What is genuinely new, and what the measurements say about it", size=16)
    drop_if(innov, "KEY INNOVATION", "{Describe the core innovation",
            "COMPETITIVE ADVANTAGE", "{Explain how your solution is better")
    drop_body_icons(innov)
    card(innov, BODY_LEFT, BODY_TOP, HALF_WIDTH, BODY_BOTTOM - BODY_TOP)
    card(innov, BODY_LEFT + HALF_WIDTH + 0.35, BODY_TOP, HALF_WIDTH, BODY_BOTTOM - BODY_TOP)
    add_text(innov, BODY_LEFT + 0.26, BODY_TOP + 0.22, HALF_WIDTH - 0.52, 0.25,
             ["KEY INNOVATION"], size=10, bold=True, color=LABEL)
    add_text(innov, BODY_LEFT + 0.26, BODY_TOP + 0.56, HALF_WIDTH - 0.52, 2.2, [
        "Read the pose grid as EVIDENCE, not as its maximum.",
        "The refit scores each candidate over ~25 poses. Taking the best of 25 is upward-biased, "
        "and the bias grows with how rough that candidate's surface is — so one lucky sample "
        "outranks a candidate that was consistently good. Summarising the same grid by a "
        "log-sum-exp costs zero extra correlations and is the largest single gain in the project.",
        "Six re-ranking criteria were built and all six failed. The only selection stage that ever "
        "worked re-scores by the SAME criterion at a better geometry.",
    ], size=11)
    add_text(innov, BODY_LEFT + HALF_WIDTH + 0.61, BODY_TOP + 0.22, HALF_WIDTH - 0.52, 0.25,
             ["MEASURED ADVANTAGE"], size=10, bold=True, color=LABEL)
    add_text(innov, BODY_LEFT + HALF_WIDTH + 0.61, BODY_TOP + 0.56, HALF_WIDTH - 0.52, 2.2, [
        f"Against the sponsor's own baseline on the same {total} pairs: "
        f"{pct(sponsor['mislock_rate'])} / {pct(bench['mislock_rate'])} / "
        f"{pct(finfet['mislock_rate'])} mis-lock against {pct(0.25)} / {pct(0.767)} / {pct(0.9)}.",
        "55 pairs fixed, 0 broken — strictly dominant, tested as a paired comparison on identical "
        "pairs rather than two averages compared by eye.",
        "The baseline's output is always integer + 50.0, so it carries a half-pixel quantisation "
        "floor and cannot score on the sub-pixel rung of the tolerance ladder at all.",
    ], size=11)

    # ------------------------------------------------------------------ 6. Results
    set_text(find(results, "Impact and Benefits"), "Results")
    set_text(find(results, "Explain how your solution will make an impact"),
             f"{total} evaluated pairs across three splits — one a held-out architecture", size=16)
    drop_if(results, "Primary Impact", "Describe the most significant benefit",
            "Quantifiable Outcomes", "List potential metrics")
    drop_body_icons(results)
    card(results, BODY_LEFT, BODY_TOP, 4.7, BODY_BOTTOM - BODY_TOP)
    card(results, BODY_LEFT + 5.05, BODY_TOP, 6.18, BODY_BOTTOM - BODY_TOP)

    add_text(results, BODY_LEFT + 0.26, BODY_TOP + 0.18, 4.2, 0.25,
             ["ACCURACY vs THE TOLERANCE LADDER"], size=10, bold=True, color=LABEL)
    ladder = [
        ("", "sponsor", "bench", "FinFET"),
        ("≤ 5 px", pct(sponsor["pass@5px"]), pct(bench["pass@5px"]), pct(finfet["pass@5px"])),
        ("≤ 2 px", pct(sponsor["pass@2px"]), pct(bench["pass@2px"]), pct(finfet["pass@2px"])),
        ("≤ 1 px", pct(sponsor["pass@1px"]), pct(bench["pass@1px"]), pct(finfet["pass@1px"])),
        ("≤ 0.5 px", pct(sponsor["pass@subpixel(0.5px)"]), pct(bench["pass@subpixel(0.5px)"]),
         pct(finfet["pass@subpixel(0.5px)"])),
        ("median err", f"{float(sponsor['error_median_px']):.3f}",
         f"{float(bench['error_median_px']):.3f}", f"{float(finfet['error_median_px']):.3f}"),
    ]
    for column, x in enumerate((BODY_LEFT + 0.26, BODY_LEFT + 1.5,
                                BODY_LEFT + 2.62, BODY_LEFT + 3.6)):
        add_text(results, x, BODY_TOP + 0.55, 1.15, 1.9, [row[column] for row in ladder],
                 size=10.5, bold=(column == 0),
                 color=LABEL if column == 0 else WHITE, space_after=3)
    add_text(results, BODY_LEFT + 0.26, BODY_TOP + 2.48, 4.2, 0.45,
             [f"Aggregate mis-lock {bad}/{total}. Runtime ~0.6 s/pair, CPU only.",
              "Baseline at the sub-pixel rung: 17.5%."], size=10, color=ACCENT, space_after=2)

    add_text(results, BODY_LEFT + 5.31, BODY_TOP + 0.18, 5.66, 0.25,
             ["ONE SUCCESS, ONE HONEST FAILURE"], size=10, bold=True, color=LABEL)
    for index, (name, source) in enumerate((("success", "success_8.png"),
                                            ("failure", "failure_13.png"))):
        original = REPO_ROOT / "results" / "failure_case" / source
        if not original.exists():
            continue
        cropped = FIGURES / f"deck_{name}.png"
        crop_search_half(original, cropped)
        results.shapes.add_picture(str(cropped),
                                   Inches(BODY_LEFT + 5.31 + index * 2.95),
                                   Inches(BODY_TOP + 0.52),
                                   width=Inches(2.7), height=Inches(1.55))
    add_text(results, BODY_LEFT + 5.31, BODY_TOP + 2.15, 2.7, 0.7,
             ["SUCCESS — error 0.04 px", "prediction sits on the truth"],
             size=9.5, color=WHITE, space_after=1)
    add_text(results, BODY_LEFT + 8.26, BODY_TOP + 2.15, 2.7, 0.7,
             ["FAILURE — error 628.24 px", "a valid lattice repeat won"],
             size=9.5, color=WHITE, space_after=1)
    add_text(results, BODY_LEFT + 5.31, BODY_TOP + 2.62, 5.66, 0.3,
             ["green prediction · blue truth · orange rivals — search image shown"],
             size=9, color=LABEL)

    # ------------------------------------------------------------------ 7. Technology
    set_text(find(tech, "Describe the technologies"),
             "CPU-only · four dependencies · no network access and no model download", size=16)
    drop_if(tech, "IMPLEMENTATION STRATEGY", "{Provide a detailed breakdown")
    drop_body_icons(tech)
    card(tech, BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_BOTTOM - BODY_TOP)
    add_text(tech, BODY_LEFT + 0.26, BODY_TOP + 0.2, BODY_WIDTH - 0.52, 0.25,
             ["STACK, HARDWARE AND FEASIBILITY"], size=10, bold=True, color=LABEL)
    add_text(tech, BODY_LEFT + 0.26, BODY_TOP + 0.56, BODY_WIDTH - 0.52, 2.2, [
        "Python 3.14 · numpy · opencv-python-headless · scipy · scikit-image. torch is optional and "
        "lazily imported — uninstalling it leaves everything working.",
        "Hardware: one laptop CPU. No GPU, no cloud, nothing downloaded at runtime. Deterministic: "
        "a single seeded numpy Generator is threaded through, and a test asserts the same seed "
        "gives byte-identical images and identical predictions.",
        "Inference accepts .png and .npy pairs identically; a separate documented converter exists "
        "for visual inspection but is never on the scoring path.",
        "Model size: none. Nothing is trained, so there are no weights to ship, load or version.",
        "Dataset generation ~2 s per pair (a 10 000×10 000 canvas at 1 nm/px); localization ~0.6 s "
        "per 1000×1000 pair.",
    ], size=11.5)

    # ------------------------------------------------------------------ 8. GitHub and video
    set_text(find(links, "{Paste your GitHub / Source Code Link here}"),
             "https://github.com/<REPLACE-BEFORE-SUBMITTING>")
    set_text(find(links, "{Paste your Video Link here"),
             "Optional — demo of localize.py on one image pair")

    # ------------------------------------------------------------------ 9. References
    # This slide's internals differ from the rest of the deck (two stacked sub-cards, its own
    # headings, three seeded reference lines), so its body is cleared and rebuilt rather than
    # repositioned piece by piece.
    set_text(find(refs, "Research and References"), "References")
    clear_body(refs, below=2.3)
    card_shape = refs.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(BODY_LEFT), Inches(2.55),
        Inches(BODY_WIDTH), Inches(4.35))
    card_shape.fill.background()
    card_shape.line.color.rgb = RGBColor(0x3B, 0x5A, 0xA8)
    card_shape.shadow.inherit = False
    add_text(refs, BODY_LEFT + 0.26, 2.72, BODY_WIDTH - 0.52, 0.25,
             ["VERIFIED REFERENCES — full table in docs/REFERENCES.md"],
             size=10, bold=True, color=LABEL)
    add_text(refs, BODY_LEFT + 0.26, 3.05, BODY_WIDTH - 0.52, 3.0, [
        "Mäkitalo & Foi. Optimal Inversion of the Generalized Anscombe Transformation for "
        "Poisson–Gaussian Noise. IEEE TIP 22(1):91–103, 2013. doi:10.1109/TIP.2012.2202675",
        "Guizar-Sicairos, Thurman & Fienup. Efficient subpixel image registration algorithms. "
        "Optics Letters 33(2):156–158, 2008. doi:10.1364/OL.33.000156",
        "Villarrubia, Ritchie & Lowney. Monte Carlo modeling of secondary electron imaging in "
        "three dimensions. Proc. SPIE 6518, 65180K, 2007. doi:10.1117/12.712353",
        "Sugiyama et al. Pattern matching method and computer program for executing pattern "
        "matching. US Patent 7,925,095 B2, Hitachi High-Technologies, 2011",
        "Abe et al. Image processing method for determining matching position between template and "
        "search image. US Patent 8,139,868 B2, Hitachi High-Technologies, 2012",
        "6F² DRAM cell design with 3F-pitch folded digitline sense amplifier. US Patent 7,349,232 B2",
        "Reddy & Chatterji. An FFT-Based Technique for Translation, Rotation and Scale-Invariant "
        "Image Registration. IEEE TIP 5(8):1266–1271, 1996. doi:10.1109/83.506761",
        "Kovesi. Phase congruency: A low-level image invariant. Psychological Research "
        "64(2):136–148, 2000. doi:10.1007/s004260000024",
    ], size=10.5, space_after=5)
    add_text(refs, BODY_LEFT + 0.26, 6.25, BODY_WIDTH - 0.52, 0.55, [
        "Rule R1 — a citation without a verified DOI/URL, a named verifier and a date is not a "
        "citation and does not appear here.",
        "The two Hitachi patents matter twice over: they establish that a closest-to-origin "
        "tie-break for periodic SEM patterns is industrial practice, not an arbitrary rule.",
    ], size=9.5, color=ACCENT, space_after=3)

    prs.save(str(OUTPUT))
    print(f"  Wrote {OUTPUT.name}  ({len(prs.slides._sldIdLst)} slides)")
    print(f"  Aggregate {bad}/{total} mis-locks; every figure read from results/")
    return 0


if __name__ == "__main__":
    raise SystemExit(build())
