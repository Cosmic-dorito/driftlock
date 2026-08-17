#!/usr/bin/env python3
"""Automated submission-readiness check.

Walks the 15-item checklist from section 9 of the Applied Materials problem statement
(see docs/SPEC.md) plus the project's own correctness rules, and reports PASS / FAIL / PENDING
for each. Run it before every push and as the last action before submitting.

    python scripts/verify_submission.py
    python scripts/verify_submission.py --strict    # PENDING counts as failure

Exit codes: 0 = no failures, 1 = at least one FAIL.

PENDING means "this deliverable does not exist yet" — expected during development, not acceptable
at submission time. Use --strict on the final run.
"""

from __future__ import annotations

import argparse
import csv
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]

# Directories that are not ours to police.
SKIP_DIRS = {".git", ".venv", "venv", "__pycache__", ".pytest_cache", ".ruff_cache",
             "third_party", "dist", "node_modules"}

PASS, FAIL, PENDING = "PASS", "FAIL", "PENDING"


@dataclass
class Result:
    status: str
    detail: str = ""
    items: list[str] = field(default_factory=list)


def iter_source_files(*suffixes: str):
    """Yield repo files with the given suffixes, skipping vendored and generated trees."""
    for path in REPO_ROOT.rglob("*"):
        if not path.is_file() or path.suffix not in suffixes:
            continue
        if any(part in SKIP_DIRS for part in path.relative_to(REPO_ROOT).parts):
            continue
        yield path


# --------------------------------------------------------------------------------------
# Checks
# --------------------------------------------------------------------------------------

def check_no_absolute_paths() -> Result:
    """Spec checklist: 'No proprietary data or hard-coded local paths are present.'

    Patterns are assembled from fragments so this file does not match itself.
    """
    drive = re.compile(r"[A-Za-z]:" + re.escape("\\") + r"[A-Za-z0-9_]")
    posix_home = re.compile(r"/(?:home|Users)/[A-Za-z0-9_.-]+/")
    offenders: list[str] = []

    for path in iter_source_files(".py", ".sh", ".yaml", ".yml", ".toml", ".cfg", ".csv"):
        if path.name == Path(__file__).name:
            continue  # this file necessarily contains the patterns it hunts for
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
        except OSError:
            continue
        for lineno, line in enumerate(text.splitlines(), 1):
            if drive.search(line) or posix_home.search(line):
                rel = path.relative_to(REPO_ROOT).as_posix()
                offenders.append(f"{rel}:{lineno}: {line.strip()[:90]}")

    if offenders:
        return Result(FAIL, f"{len(offenders)} hard-coded path(s)", offenders[:15])
    return Result(PASS, "no hard-coded absolute paths in source")


def check_no_os_path() -> Result:
    """Portability: pathlib only. os.path assumptions break across Windows and Linux."""
    offenders = []
    for path in iter_source_files(".py"):
        if path.name == Path(__file__).name:
            continue
        text = path.read_text(encoding="utf-8", errors="ignore")
        for lineno, line in enumerate(text.splitlines(), 1):
            if re.search(r"\bos\.path\.", line) or re.search(r"\bos\.sep\b", line):
                offenders.append(f"{path.relative_to(REPO_ROOT).as_posix()}:{lineno}")
    if offenders:
        return Result(FAIL, f"{len(offenders)} os.path use(s); use pathlib", offenders[:10])
    return Result(PASS, "pathlib only")


def check_required_files() -> Result:
    """Spec deliverables 2, 3, 4: separate runnable generator + localizer, README, dependencies."""
    required = ["README.md", "requirements.txt", "generate_dataset.py", "localize.py", "LICENSE"]
    missing = [f for f in required if not (REPO_ROOT / f).exists()]
    if missing:
        return Result(PENDING, f"not yet created: {', '.join(missing)}")
    return Result(PASS, "all required top-level deliverables present")


def check_generator_and_localizer_separate() -> Result:
    """Spec checklist: 'Python generator and Python localization code are separate and runnable.'"""
    gen, loc = REPO_ROOT / "generate_dataset.py", REPO_ROOT / "localize.py"
    if not (gen.exists() and loc.exists()):
        return Result(PENDING, "generator and/or localizer not yet written")
    return Result(PASS, "generate_dataset.py and localize.py are separate entry points")


def check_no_notebook_only() -> Result:
    """Spec: 'A single notebook is not sufficient as the only runnable submission.'"""
    scripts = [p for p in iter_source_files(".py") if p.parent == REPO_ROOT]
    if not scripts:
        return Result(PENDING, "no top-level Python entry points yet")
    return Result(PASS, f"{len(scripts)} runnable script(s) at repo root")


def check_bench_dataset() -> Result:
    """Spec checklist: 'At least 30 varied cases are evaluated.'"""
    manifest = REPO_ROOT / "data" / "bench" / "manifest.csv"
    if not manifest.exists():
        return Result(PENDING, "data/bench/manifest.csv not yet generated")
    with manifest.open(newline="", encoding="utf-8") as fh:
        rows = list(csv.DictReader(fh))
    if len(rows) < 30:
        return Result(FAIL, f"only {len(rows)} pairs; the spec requires at least 30")
    return Result(PASS, f"{len(rows)} validation pairs")


def check_manifest_columns() -> Result:
    """Spec deliverable 6: manifest carries paths, ground truth, predictions and metadata."""
    manifest = REPO_ROOT / "data" / "bench" / "manifest.csv"
    if not manifest.exists():
        return Result(PENDING, "manifest not yet generated")
    required = {"id", "reference_path", "search_path", "gt_x", "gt_y", "seed",
                "architecture", "scale_ratio", "rotation_deg"}
    with manifest.open(newline="", encoding="utf-8") as fh:
        cols = set(next(csv.reader(fh)))
    missing = required - cols
    if missing:
        return Result(FAIL, f"manifest missing columns: {sorted(missing)}")
    return Result(PASS, f"manifest has all required columns ({len(cols)} total)")


def check_image_dimensions() -> Result:
    """Spec checklist: '1000 x 1000 reference and search images used.'"""
    ref_dir = REPO_ROOT / "data" / "bench" / "reference"
    if not ref_dir.exists() or not any(ref_dir.glob("*.png")):
        return Result(PENDING, "no bench images yet")
    try:
        from PIL import Image
    except ImportError:
        return Result(PENDING, "pillow not installed; cannot check dimensions")
    bad = []
    for img_path in sorted(ref_dir.glob("*.png"))[:5]:
        with Image.open(img_path) as im:
            if im.size != (1000, 1000):
                bad.append(f"{img_path.name}: {im.size}")
    if bad:
        return Result(FAIL, "images are not 1000x1000", bad)
    return Result(PASS, "reference images are 1000x1000")


def _metrics_files() -> list[Path]:
    """Every metrics CSV `evaluate.py` may have written.

    It emits one per split as ``metrics_<label>.csv`` (and plain ``metrics.csv`` when unlabelled),
    because reporting a single blended number across splits with different difficulty would be the
    kind of average that hides the failure mode (R9). The check therefore globs rather than
    insisting on one filename.
    """
    results = REPO_ROOT / "results"
    return sorted(results.glob("metrics*.csv")) if results.is_dir() else []


def check_metrics_reported() -> Result:
    """Spec checklist: '5-, 4-, 2- and 1-pixel pass rates are reported.'"""
    files = _metrics_files()
    if not files:
        return Result(PENDING, "no results/metrics*.csv generated yet")
    needed = ["pass@5", "pass@4", "pass@2", "pass@1"]
    for metrics in files:
        text = metrics.read_text(encoding="utf-8").lower()
        missing = [n for n in needed if n not in text]
        if missing:
            return Result(FAIL, f"{metrics.name} missing: {missing}")
    return Result(PASS, f"threshold-wise pass rates reported in {len(files)} metrics file(s)")


def check_runtime_reported() -> Result:
    """Spec checklist: 'Runtime, hardware and timing method are stated.'"""
    files = _metrics_files()
    if not files:
        return Result(PENDING, "results not yet generated")
    text = files[0].read_text(encoding="utf-8").lower()
    if "runtime" not in text:
        return Result(FAIL, f"no runtime in {files[0].name}")
    readme = (REPO_ROOT / "README.md").read_text(encoding="utf-8").lower()
    if "python version" not in readme and "hardware" not in readme:
        return Result(PENDING, "runtime present, but hardware/timing method not yet stated in README")
    return Result(PASS, "runtime, hardware and timing method stated")


def check_failure_case() -> Result:
    """Spec checklist: 'At least one failure case is shown and explained.'"""
    d = REPO_ROOT / "results" / "failure_case"
    if not d.exists() or not any(d.iterdir()):
        return Result(PENDING, "results/failure_case/ is empty")
    has_image = any(d.glob("*.png"))
    has_text = any(d.glob("*.md"))
    if not (has_image and has_text):
        return Result(FAIL, "failure case needs both a visualization (.png) and an explanation (.md)")
    return Result(PASS, "failure case visualized and explained")


def check_robustness_sweep() -> Result:
    """Spec deliverable: 'Results across multiple noise levels, target positions, scales and
    rotations.'

    This was missing for most of the project's life, which is why it is a check and not a habit.
    Every other required deliverable had a verifier line; this one did not, so nothing complained
    that all reported numbers came from a single operating point.
    """
    path = REPO_ROOT / "results" / "robustness.csv"
    if not path.exists():
        return Result(PENDING, "run scripts/robustness_sweep.py")
    with path.open(newline="", encoding="utf-8") as fh:
        rows = [r for r in csv.DictReader(fh) if r.get("axis") and not r["axis"].startswith("#")]
    axes = {r["axis"].split(" ")[0] for r in rows}
    missing = {"noise", "scale", "rotation"} - axes
    if missing:
        return Result(FAIL, f"robustness sweep is missing axes: {sorted(missing)}")
    # The spec names FOUR axes and target position is the one that needs no new data - the evaluated
    # pairs already span 32-547 px from the field centre - so it is stratified rather than swept,
    # and lives in its own file. Checked here so the deliverable cannot be three-quarters met.
    if not (REPO_ROOT / "results" / "position_strata.csv").exists():
        return Result(FAIL, "target-position stratification missing: run scripts/position_strata.py")
    return Result(PASS, f"{len(rows)} operating points across {len(axes)} axes, plus position")


def check_failure_decomposition() -> Result:
    """Not a spec item - our own. Failure analysis is 10% of the score, and a single mis-lock rate
    does not say which STAGE lost each pair. Fails if the buckets go missing or stop covering the
    reported splits, since the deck quotes them.
    """
    path = REPO_ROOT / "results" / "failure_decomposition.csv"
    if not path.exists():
        return Result(PENDING, "run scripts/failure_decomposition.py")
    with path.open(newline="", encoding="utf-8") as fh:
        rows = [r for r in csv.DictReader(fh) if r.get("split") and not r["split"].startswith("#")]
    buckets = {r["bucket"] for r in rows}
    if not {"correct", "absent", "screened", "outscored"} & buckets:
        return Result(FAIL, "failure decomposition has no recognised buckets")
    splits = {r["split"] for r in rows}
    if len(splits) < 3:
        return Result(FAIL, f"decomposition covers only {sorted(splits)}; needs all three splits")
    return Result(PASS, f"{len(rows)} pairs classified across {len(splits)} splits")


def check_citations_complete() -> Result:
    """Rule R1 + spec checklist: 'At least 2-3 credible public sources are cited.'

    Every row in the Verified table needs all five columns filled. PARTIAL rows are flagged: they
    must be upgraded or dropped before they can appear in the deck.
    """
    refs = REPO_ROOT / "docs" / "REFERENCES.md"
    if not refs.exists():
        return Result(FAIL, "docs/REFERENCES.md missing")
    text = refs.read_text(encoding="utf-8")
    verified = len(re.findall(r"\|\s*VERIFIED\s*\|", text))
    # Count PARTIAL only where it is a row's STATUS - i.e. inside a table cell. The word also
    # appears in the file's own status legend, and counting that made the check permanently
    # unsatisfiable no matter how many citations were upgraded.
    partial = len([line for line in text.splitlines()
                   if line.lstrip().startswith("|") and "PARTIAL" in line])
    if verified < 3:
        return Result(FAIL, f"only {verified} fully VERIFIED citation(s); the spec requires 2-3 credible sources")
    detail = f"{verified} verified"
    if partial:
        return Result(PENDING, f"{detail}, but {partial} PARTIAL row(s) must be upgraded or dropped")
    return Result(PASS, detail)


def check_ppt_numbers_traceable() -> Result:
    """Rule R2: every number in the deck must exist in results/.

    The likely failure is a STALE number surviving from an earlier run, not an invented one.
    """
    decks = list(REPO_ROOT.glob("*.pptx")) + list(REPO_ROOT.glob("**/solution_presentation.pptx"))
    # Skip Office/LibreOffice lock files ("~$name.pptx"). They are not decks, they are not readable
    # while the real file is open, and treating one as a deck fails this check for a reason that has
    # nothing to do with the submission - which is exactly the kind of false alarm that teaches
    # people to ignore a verifier.
    decks = [d for d in decks
             if not d.name.startswith("~$")
             and not any(p in SKIP_DIRS for p in d.relative_to(REPO_ROOT).parts)]
    if not decks:
        return Result(PENDING, "no .pptx yet - mandatory deliverable 1")
    try:
        from pptx import Presentation
    except ImportError:
        return Result(PENDING, "python-pptx not installed (pip install -r requirements-dev.txt)")

    results_text = " ".join(
        p.read_text(encoding="utf-8", errors="ignore")
        for p in (REPO_ROOT / "results").rglob("*") if p.is_file() and p.suffix in {".csv", ".md", ".json"}
    )
    if not results_text.strip():
        return Result(FAIL, "deck exists but results/ is empty — no number can be traced")

    unmatched = []
    for deck in decks:
        for slide_no, slide in enumerate(Presentation(deck).slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for num in re.findall(r"\d+\.\d+", shape.text_frame.text):
                    if num not in results_text:
                        unmatched.append(f"{deck.name} slide {slide_no}: {num}")
    if unmatched:
        return Result(FAIL, f"{len(unmatched)} deck number(s) not found in results/", unmatched[:12])
    return Result(PASS, "every decimal number in the deck is traceable to results/")


def check_results_doc_is_current() -> Result:
    """R2 for the markdown, not just the deck.

    The deck is generated from results/ and checked number-by-number, but docs/RESULTS.md was
    hand-maintained - and it drifted, to the point where RESULTS.md, PROGRESS.md and results/ each
    carried a different generation of the headline figures at the same time. A judge browsing the
    repository would have found the project contradicting itself.

    The headline block is now generated between markers. This check confirms it is present and that
    its numbers still match results/, so a stale document fails the build instead of shipping.
    """
    doc = REPO_ROOT / "docs" / "RESULTS.md"
    if not doc.exists():
        return Result(FAIL, "docs/RESULTS.md missing")
    text = doc.read_text(encoding="utf-8")
    if "BEGIN GENERATED HEADLINE" not in text:
        return Result(FAIL, "RESULTS.md headline is not generated - run scripts/make_results_doc.py")

    block = text.split("BEGIN GENERATED HEADLINE")[1].split("END GENERATED HEADLINE")[0]
    stale = []
    for label in ("sponsor", "bench", "finfet"):
        path = REPO_ROOT / "results" / f"metrics_{label}.csv"
        if not path.exists():
            continue
        with path.open(newline="", encoding="utf-8") as fh:
            metrics = {r["metric"]: r["value"] for r in csv.DictReader(fh)}
        expected = f"{float(metrics['mislock_rate']) * 100:.1f}%"
        if expected not in block:
            stale.append(f"{label}: results/ says {expected}, not found in the generated block")
    if stale:
        return Result(FAIL, "RESULTS.md is stale - re-run scripts/make_results_doc.py", stale)
    return Result(PASS, "RESULTS.md headline is generated and matches results/")


def check_config_blocks_are_generated() -> Result:
    """Every document that states the shipped configuration must state the same one.

    Found by an external reviewer, not by us: after ADR-0037 moved the screen cut to 40,
    `localize.py` and STATE.md agreed while HANDOFF.md still said 30. Three hand-copied duplicates
    of a nine-line config drifted within hours of the change that created them. The blocks are now
    generated from ``build_config`` by scripts/make_results_doc.py; this confirms they are present
    and current, so the next drift fails the build instead of reaching a reviewer.
    """
    import argparse as _ap
    from dataclasses import fields

    if str(REPO_ROOT) not in sys.path:
        sys.path.insert(0, str(REPO_ROOT))
    try:
        import localize as L
        from src.driftlock.match import PipelineConfig
    except ImportError as exc:                      # pragma: no cover - environment problem
        return Result(PENDING, f"cannot import the pipeline to check config blocks ({exc})")

    shipped = L.build_config(_ap.Namespace(config="driftlock"))
    defaults = PipelineConfig()
    expected = [f"{f.name}={getattr(shipped, f.name)!r}" for f in fields(shipped)
                if getattr(shipped, f.name) != getattr(defaults, f.name)]

    problems = []
    for rel in ("docs/STATE.md", "docs/HANDOFF.md"):
        path = REPO_ROOT / rel
        if not path.exists():
            continue
        text = path.read_text(encoding="utf-8", errors="ignore")
        if "BEGIN GENERATED CONFIG" not in text:
            problems.append(f"{rel}: config block is not generated")
            continue
        block = text.split("BEGIN GENERATED CONFIG")[1].split("END GENERATED CONFIG")[0]
        missing = [item for item in expected if item not in block]
        if missing:
            problems.append(f"{rel}: stale, missing {', '.join(missing[:4])}")
    if problems:
        return Result(FAIL, "config blocks disagree with build_config - "
                            "run scripts/make_results_doc.py", problems)
    return Result(PASS, f"config blocks match build_config ({len(expected)} non-default fields)")


def check_results_newer_than_config() -> Result:
    """Every measured artefact must post-date the configuration it claims to describe.

    `check_results_doc_is_current` compares documents against `results/`. Nothing compared
    `results/` against the *code*, and that is a real gap rather than a hypothetical one:
    `results/robustness.csv` was last regenerated at 68d1ef0 and survived the ADR-0035 config change
    untouched, so its 25 operating points described a pipeline that no longer existed. Four
    documents quoted it. Nothing failed, because every number in it was internally consistent - it
    was consistently describing the wrong build.

    Mtime is a weak signal, so this only WARNS via PENDING rather than failing the build: a fresh
    clone has arbitrary timestamps, and a legitimate comment-only edit to localize.py would trip it.
    It is here to be noticed at submission time, not to gate a checkout.
    """
    config_files = [REPO_ROOT / "localize.py",
                    REPO_ROOT / "src" / "driftlock" / "match.py",
                    REPO_ROOT / "src" / "driftlock" / "refit.py"]
    config_files = [p for p in config_files if p.exists()]
    if not config_files:
        return Result(PENDING, "no configuration files found to compare against")
    newest_config = max(p.stat().st_mtime for p in config_files)

    # Only artefacts produced by RUNNING the pipeline. Derived documents are covered elsewhere.
    measured = ["robustness.csv", "significance.csv", "failure_decomposition.csv",
                "metrics_sponsor.csv", "metrics_bench.csv", "metrics_finfet.csv", "ablation.md"]
    stale = []
    for name in measured:
        path = REPO_ROOT / "results" / name
        if not path.exists():
            continue
        if path.stat().st_mtime < newest_config:
            stale.append(f"{name} predates the current pipeline configuration")
    if stale:
        return Result(PENDING, f"{len(stale)} measured artefact(s) may predate the config", stale)
    return Result(PASS, "every measured artefact post-dates the current configuration")


def check_torch_optional() -> Result:
    """ADR-0006: the deterministic path must not import torch at module scope."""
    offenders = []
    for path in iter_source_files(".py"):
        if "rerank" in path.name or path.parts[-2:][0] == "tests":
            continue
        for lineno, line in enumerate(path.read_text(encoding="utf-8", errors="ignore").splitlines(), 1):
            stripped = line.strip()
            if re.match(r"^(import torch|from torch)", stripped):
                indent = len(line) - len(line.lstrip())
                if indent == 0:  # module scope, not inside a function
                    offenders.append(f"{path.relative_to(REPO_ROOT).as_posix()}:{lineno}")
    if offenders:
        return Result(FAIL, "torch imported at module scope; it must be lazy", offenders)
    return Result(PASS, "torch is not imported at module scope")


def check_docs_present() -> Result:
    """Project practice: the tracking docs that make the repo resumable on another machine."""
    required = ["PLAN.md", "SPEC.md", "DECISIONS.md", "PROGRESS.md", "HANDOFF.md", "REFERENCES.md"]
    missing = [f for f in required if not (REPO_ROOT / "docs" / f).exists()]
    if missing:
        return Result(FAIL, f"docs missing: {', '.join(missing)}")
    return Result(PASS, "all tracking docs present")


def check_hypotheses_verified() -> Result:
    """Rule R3: the sponsor-generator facts were read from source, not run. Confirm before relying."""
    decisions = REPO_ROOT / "docs" / "DECISIONS.md"
    if not decisions.exists():
        return Result(FAIL, "docs/DECISIONS.md missing")
    unverified = len(re.findall(r"\|\s*unverified\s*\|", decisions.read_text(encoding="utf-8")))
    if unverified:
        return Result(PENDING, f"{unverified} of H1-H9 still unverified - B owns this on Day 1")
    return Result(PASS, "all sponsor-generator hypotheses empirically confirmed")


CHECKS = [
    ("No hard-coded local paths", check_no_absolute_paths),
    ("pathlib only (no os.path)", check_no_os_path),
    ("Required deliverable files present", check_required_files),
    ("Generator and localizer are separate", check_generator_and_localizer_separate),
    ("Runnable scripts, not notebook-only", check_no_notebook_only),
    ("At least 30 validation pairs", check_bench_dataset),
    ("Manifest has required columns", check_manifest_columns),
    ("Images are 1000x1000", check_image_dimensions),
    ("Pass rates at 5/4/2/1 px reported", check_metrics_reported),
    ("Runtime, hardware, timing method", check_runtime_reported),
    ("Failure case shown and explained", check_failure_case),
    ("Noise/scale/rotation sweep reported", check_robustness_sweep),
    ("Failures split by losing stage", check_failure_decomposition),
    ("Citations complete and verified (R1)", check_citations_complete),
    ("Deck numbers traceable to results (R2)", check_ppt_numbers_traceable),
    ("RESULTS.md generated and current (R2)", check_results_doc_is_current),
    ("Config blocks match build_config", check_config_blocks_are_generated),
    ("Measured artefacts post-date the config", check_results_newer_than_config),
    ("torch is optional, lazily imported", check_torch_optional),
    ("Tracking docs present", check_docs_present),
    ("Sponsor-generator facts verified (R3)", check_hypotheses_verified),
]


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--strict", action="store_true",
                        help="treat PENDING as failure (use for the final pre-submission run)")
    args = parser.parse_args()

    symbol = {PASS: "PASS", FAIL: "FAIL", PENDING: "PEND"}
    counts = {PASS: 0, FAIL: 0, PENDING: 0}
    width = max(len(name) for name, _ in CHECKS)

    # ASCII only in console output: Windows consoles default to cp1252 and mangle non-ASCII.
    print(f"\nDriftLock submission check  --  {REPO_ROOT}\n")
    for name, fn in CHECKS:
        try:
            res = fn()
        except Exception as exc:  # a broken check must not look like a passing one
            res = Result(FAIL, f"check raised {type(exc).__name__}: {exc}")
        counts[res.status] += 1
        print(f"  [{symbol[res.status]}] {name:<{width}}  {res.detail}")
        for item in res.items:
            print(f"         - {item}")

    print(f"\n  {counts[PASS]} passed, {counts[FAIL]} failed, {counts[PENDING]} pending\n")

    if counts[FAIL]:
        print("  FAILED. Fix the items above before pushing.\n")
        return 1
    if counts[PENDING] and args.strict:
        print("  STRICT MODE: pending items are not acceptable at submission time.\n")
        return 1
    if counts[PENDING]:
        print("  No failures. Pending items are deliverables not yet built - expected during\n"
              "  development. Re-run with --strict before submitting.\n")
    else:
        print("  All checks passed.\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
